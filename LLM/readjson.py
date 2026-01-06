import os
import json
import pdfplumber
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation

# ========== PDF ==========
def pdf_to_json(file_path):
    data = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text()
            data.append({
                "page": i,
                "content": text if text else ""
            })
    return data

# ========== Word (.docx) ==========
def docx_to_json(file_path):
    doc = Document(file_path)
    data = []
    for i, para in enumerate(doc.paragraphs, start=1):
        if para.text.strip():
            data.append({
                "paragraph": i,
                "content": para.text.strip()
            })
    return data

# ========== PowerPoint (.pptx) ==========
def pptx_to_json(file_path):
    prs = Presentation(file_path)
    data = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        data.append({
            "slide": i,
            "content": slide_text
        })
    return data

# ========== Image (jpg/png) ==========
def image_to_json(file_path, lang="eng"):
    img = Image.open(file_path)
    text = pytesseract.image_to_string(img, lang=lang)
    return [{"image": os.path.basename(file_path), "content": text.strip()}]

# ========== Main ==========
def convert_file(file_path, lang="eng"):
    ext = os.path.splitext(file_path)[-1].lower()
    if ext == ".pdf":
        return pdf_to_json(file_path)
    elif ext == ".docx":
        return docx_to_json(file_path)
    elif ext == ".pptx":
        return pptx_to_json(file_path)
    elif ext in [".jpg", ".jpeg", ".png"]:
        return image_to_json(file_path, lang=lang)
    else:
        raise ValueError("Unsupported file type")

if __name__ == "__main__":
    # Thay đường dẫn file cần convert
    file_path = "C:/Users/lengu/Desktop/Tổng hợp nội dung bài báo.docx"  # có thể đổi sang .docx, .pptx, .jpg
    
    result = convert_file(file_path, lang="eng")
    
    # Xuất ra JSON
    out_file = os.path.splitext(file_path)[0] + ".json"
    with open(out_file, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    
    print(f"✅ Chuyển đổi thành công! File JSON lưu tại: {out_file}")
